#!/usr/bin/env python3
"""Insert generated PPT page images into a PowerPoint wrapper.

Requires python-pptx and Pillow:
    python -m pip install python-pptx pillow
"""

from __future__ import annotations

import argparse
import re
from pathlib import Path


def slide_sort_key(path: Path) -> tuple[int, str]:
    match = re.search(r"slide[-_ ]?(\d+)", path.stem, re.IGNORECASE)
    if match:
        return (int(match.group(1)), path.name)
    return (10_000, path.name)


def import_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ModuleNotFoundError as exc:
        raise SystemExit(
            "Missing dependency: python-pptx. Install it with:\n"
            "  python -m pip install python-pptx"
        ) from exc
    return Presentation, Inches


def import_pillow():
    try:
        from PIL import Image
    except ModuleNotFoundError as exc:
        raise SystemExit(
            "Missing dependency: Pillow. Install it with:\n"
            "  python -m pip install pillow"
        ) from exc
    return Image


def add_image_page(slide, image_path: Path, slide_width: int, slide_height: int) -> None:
    Image = import_pillow()
    with Image.open(image_path) as img:
        image_width, image_height = img.size

    if image_width <= 0 or image_height <= 0:
        raise SystemExit(f"Invalid image size for {image_path}")

    slide_aspect = slide_width / slide_height
    image_aspect = image_width / image_height

    if image_aspect >= slide_aspect:
        pic_width = slide_width
        pic_height = int(slide_width / image_aspect)
        left = 0
        top = int((slide_height - pic_height) / 2)
    else:
        pic_height = slide_height
        pic_width = int(slide_height * image_aspect)
        left = int((slide_width - pic_width) / 2)
        top = 0

    slide.shapes.add_picture(str(image_path), left, top, width=pic_width, height=pic_height)


def build_pptx(image_dir: Path, output: Path, pattern: str) -> None:
    Presentation, Inches = import_pptx()
    files = sorted(image_dir.glob(pattern), key=slide_sort_key)
    if not files:
        raise SystemExit(f"No images matched {pattern!r} in {image_dir}")

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for file in files:
        slide = prs.slides.add_slide(blank_layout)
        add_image_page(slide, file, prs.slide_width, prs.slide_height)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("image_dir", type=Path, nargs="?", help="Directory containing final PPT page images")
    parser.add_argument("--input-dir", type=Path, default=None, help="Directory containing final PPT page images")
    parser.add_argument("-o", "--output", type=Path, required=True, help="Output .pptx path")
    parser.add_argument("--pattern", default="slide-*.png", help="Glob pattern, default: slide-*.png")
    args = parser.parse_args()

    image_dir = args.input_dir or args.image_dir
    if image_dir is None:
        parser.error("provide an image directory as a positional argument or with --input-dir")

    build_pptx(image_dir, args.output, args.pattern)
    print(args.output)


if __name__ == "__main__":
    main()
